# -*- coding: utf-8 -*-
"""꿈제작소 AR 게임 제안서 PPTX 생성"""
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn as _qn
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

# ===== 색상 =====
TEAL_DARK = RGBColor(0x1D, 0x45, 0x56)
TEAL      = RGBColor(0x32, 0xE5, 0xC1)
TEAL_SOFT = RGBColor(0xE8, 0xF7, 0xF4)
IVORY     = RGBColor(0xF7, 0xF1, 0xE3)
WOOD      = RGBColor(0xC9, 0x90, 0x5F)
YELLOW    = RGBColor(0xFD, 0xF0, 0xB5)
YELLOW_D  = RGBColor(0xE8, 0xC5, 0x4A)
GRAY_TXT  = RGBColor(0x44, 0x44, 0x44)
GRAY_MID  = RGBColor(0x77, 0x77, 0x77)
BOX_BG    = RGBColor(0xF8, 0xF8, 0xF8)
BORDER    = RGBColor(0xDF, 0xDF, 0xDF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
RED_SOFT  = RGBColor(0xC0, 0x39, 0x2B)
GREEN_S   = RGBColor(0x27, 0xAE, 0x60)

FONT = '맑은 고딕'

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW = 13.333
SH = 7.5


# ===== 유틸 =====
def set_ea_font(font_obj, name=FONT):
    """한글(East Asian) 폰트 지정"""
    font_obj.name = name
    rPr = font_obj._element
    for tag in ('a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set('typeface', name)


def style_para(p, size=10, bold=False, color=GRAY_TXT, align=PP_ALIGN.LEFT,
               space_after=2, line=1.35):
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        set_ea_font(r.font)
    if not p.runs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        set_ea_font(p.font)


def add_tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    return tf


def add_text(slide, x, y, w, h, text, size=10, bold=False, color=GRAY_TXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=1.35):
    tf = add_tb(slide, x, y, w, h, anchor)
    lines = text.split('\n')
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        style_para(p, size=size, bold=bold, color=color, align=align, line=line)
    return tf


def rounded(slide, x, y, w, h, fill=BOX_BG, line_color=BORDER,
            line_w=1.0, radius=0.06):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    sh.text_frame.text = ''
    return sh


def rect(slide, x, y, w, h, fill, line_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    sh.text_frame.text = ''
    return sh


def oval(slide, cx, cy, d, fill=YELLOW, line_color=YELLOW_D):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(cx - d / 2), Inches(cy - d / 2),
                                Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1.2)
    sh.shadow.inherit = False
    return sh


def arrow(slide, x1, y1, x2, y2, color=RGBColor(0xE8, 0x4C, 0x3D), width=2.0):
    """화살촉이 있는 직선 화살표"""
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    ln = cn.line._get_or_add_ln()
    for tag in ('a:tailEnd',):
        el = ln.find(_qn(tag))
        if el is None:
            el = etree.SubElement(ln, _qn(tag))
        el.set('type', 'triangle')
        el.set('w', 'med')
        el.set('len', 'med')
    return cn


def edge_arrow(slide, c1, c2, d, gap=0.11, **kw):
    """두 원의 테두리 사이를 잇는 화살표"""
    (ax, ay), (bx, by) = c1, c2
    dx, dy = bx - ax, by - ay
    dist = math.hypot(dx, dy)
    ux, uy = dx / dist, dy / dist
    off = d / 2 + gap
    return arrow(slide, ax + ux * off, ay + uy * off,
                 bx - ux * off, by - uy * off, **kw)


def box(slide, x, y, w, h, title=None, body=None, fill=BOX_BG,
        line_color=BORDER, title_size=13, body_size=10,
        title_color=TEAL_DARK, body_color=GRAY_TXT,
        align=PP_ALIGN.CENTER, radius=0.06, title_bold=True):
    """제목 + 본문 박스"""
    rounded(slide, x, y, w, h, fill=fill, line_color=line_color, radius=radius)
    pad = 0.22
    tf = add_tb(slide, x + pad, y + pad * 0.75, w - pad * 2, h - pad * 1.5,
                anchor=MSO_ANCHOR.TOP)
    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        style_para(p, size=title_size, bold=title_bold, color=title_color,
                   align=align, space_after=6)
        first = False
    if body:
        for ln in body.split('\n'):
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.text = ln
            style_para(p, size=body_size, color=body_color, align=align,
                       space_after=3, line=1.4)
    return tf


PHONE_W = 2.08   # 세로형 휴대폰 스크린샷 기준 폭
PHONE_H = 4.5    # 높이 (비율 약 0.46 : 실제 스크린샷과 일치)


def phone_slot(slide, x, y, caption, hint='', w=PHONE_W, h=PHONE_H):
    """세로형 휴대폰 스크린샷을 넣을 자리. 캡션이 아래에 붙는다."""
    sh = rounded(slide, x, y, w, h, fill=RGBColor(0xF3, 0xF3, 0xF3),
                 line_color=RGBColor(0xC8, 0xC8, 0xC8), radius=0.04)
    # 안내 문구
    tf = add_tb(slide, x + 0.1, y, w - 0.2, h, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.text = '스크린샷'
    style_para(p, size=9, bold=True, color=RGBColor(0x99, 0x99, 0x99),
               align=PP_ALIGN.CENTER, space_after=4)
    p2 = tf.add_paragraph()
    p2.text = hint or caption
    style_para(p2, size=8.5, color=RGBColor(0xAA, 0xAA, 0xAA),
               align=PP_ALIGN.CENTER, line=1.35)
    # 캡션
    add_text(slide, x, y + h + 0.08, w, 0.3, caption, size=9, bold=True,
             color=TEAL_DARK, align=PP_ALIGN.CENTER)
    return sh


def slide_header(slide, num, title, sub=None):
    """상단 섹션 번호 + 제목 + 구분선. num 이 비어 있으면 번호 원을 생략한다."""
    if num:
        d = 0.62
        cx, cy = 1.42, 0.86
        oval(slide, cx, cy, d, fill=YELLOW, line_color=None)
        tf = add_tb(slide, cx - d / 2, cy - d / 2 + 0.04, d, d,
                    anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.text = num
        style_para(p, size=12, bold=True, color=TEAL_DARK,
                   align=PP_ALIGN.CENTER)
        title_x = 1.92
    else:
        title_x = 0.95
    # 제목
    add_text(slide, title_x, 0.44, 10.6, 0.9, title, size=30, bold=True,
             color=TEAL_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # 구분선
    rect(slide, 0.9, 1.42, SW - 1.8, 0.02, RGBColor(0xCC, 0xCC, 0xCC))
    if sub:
        add_text(slide, 0.92, 1.52, SW - 1.84, 0.3, sub, size=10.5,
                 color=GRAY_MID, align=PP_ALIGN.LEFT)


def footer(slide, page):
    rect(slide, 0.9, 6.94, SW - 1.8, 0.015, RGBColor(0xEA, 0xEA, 0xEA))
    add_text(slide, 0.9, 7.0, 6.0, 0.3, '꿈제작소 AR 게임 · YOUTH ECHO',
             size=8, color=RGBColor(0xAA, 0xAA, 0xAA))
    add_text(slide, SW - 2.4, 7.0, 1.5, 0.3, str(page), size=9,
             color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)


_page_no = [1]  # 표지가 1페이지


def new_slide(num=None, title=None, sub=None, page=None):
    """페이지 번호는 자동으로 매겨진다. page 인자는 하위 호환용으로 무시한다."""
    _page_no[0] += 1
    s = prs.slides.add_slide(BLANK)
    if title:
        slide_header(s, num, title, sub)
    footer(s, _page_no[0])
    return s


# =====================================================================
# 1. 표지
# =====================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, TEAL_DARK)
# 장식
rect(s, 0, 0, SW, 0.14, TEAL)
oval(s, SW - 1.2, 0.9, 5.6, fill=RGBColor(0x24, 0x52, 0x64), line_color=None)
oval(s, 0.6, SH - 0.5, 4.2, fill=RGBColor(0x23, 0x4E, 0x60), line_color=None)

add_text(s, 1.5, 1.95, 10, 0.4, 'D R E A M   W O R K S H O P',
         size=12.5, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
add_text(s, 1.45, 2.45, 11, 1.5, '꿈제작소 AR 게임', size=54, bold=True,
         color=WHITE, align=PP_ALIGN.LEFT)
rect(s, 1.5, 3.98, 1.5, 0.05, TEAL)
add_text(s, 1.5, 4.22, 10.5, 0.9,
         'YOUTH ECHO  ·  청년의 날 현장 체험 콘텐츠 제안',
         size=17, color=RGBColor(0xD8, 0xE4, 0xE8), align=PP_ALIGN.LEFT)
add_text(s, 1.5, 5.0, 10.5, 0.9,
         '꿈제작소에서만 플레이할 수 있는 위치 기반 모바일 AR 게임',
         size=11.5, color=RGBColor(0x9F, 0xB8, 0xC0), align=PP_ALIGN.LEFT)

rect(s, 1.5, 6.15, 4.4, 0.02, RGBColor(0x3E, 0x6B, 0x7C))
add_text(s, 1.5, 6.3, 6.5, 0.6,
         '기획 · 개발   이 아 선\npenguee@naver.com',
         size=10.5, color=RGBColor(0x8F, 0xAC, 0xB6), align=PP_ALIGN.LEFT)

# =====================================================================
# 2. 목차
# =====================================================================
s = new_slide('', '목차')

toc = [
    ('01', '기획 의도', '어떤 계기로 만들었고, 무엇을 기대할 수 있는가', 3),
    ('02', '게임의 구조', 'AR · 도감 · 스토리 · 배틀의 네 가지 축', 6),
    ('03', '화면 소개', '실제 플레이 화면과 게임 규칙', 7),
    ('04', '기관 운영 안내', '개인정보 · 준비물 · 행사 부스 운영', 17),
    ('05', '활용 흐름', '홍보부터 경품 수령까지의 동선', 20),
    ('06', '진행 제안', '완성 범위와 결정이 필요한 항목', 21),
]
y = 1.8
for num, title, desc, pg in toc:
    rounded(s, 0.9, y, 11.5, 0.64, fill=BOX_BG, line_color=None, radius=0.06)
    # 번호
    add_text(s, 1.25, y + 0.1, 0.9, 0.45, num, size=17, bold=True,
             color=RGBColor(0xC9, 0xD4, 0xD9), align=PP_ALIGN.LEFT)
    # 제목
    add_text(s, 2.25, y + 0.06, 3.2, 0.32, title, size=13.5, bold=True,
             color=TEAL_DARK, align=PP_ALIGN.LEFT)
    # 설명
    add_text(s, 2.25, y + 0.36, 7.8, 0.26, desc, size=9, color=GRAY_MID,
             align=PP_ALIGN.LEFT)
    # 페이지
    add_text(s, 11.05, y + 0.14, 1.05, 0.36, str(pg), size=14, bold=True,
             color=TEAL_DARK, align=PP_ALIGN.RIGHT)
    y += 0.72

box(s, 0.9, 6.2, 11.5, 0.6,
    body='이 게임은 꿈제작소 현장에서만 플레이할 수 있으며, 참가자 개인정보를 수집하지 않습니다.',
    fill=IVORY, line_color=WOOD, body_size=10.5)

# =====================================================================
# 3. 01 기획 의도 — 인사말
# =====================================================================
s = new_slide('01', 'AR게임 기획 의도', page=2)

box(s, 0.9, 1.85, 7.1, 4.8, fill=BOX_BG, radius=0.03)
tf = add_tb(s, 1.25, 2.15, 6.4, 4.2)
lines = [
    ('안녕하세요.', 14, True, TEAL_DARK, 12),
    ('내일꿈제작소 서포터즈 이아선입니다.', 11, False, GRAY_TXT, 4),
    ("면접에서 '청년의 날' 행사 아이디어를 말씀하셨던 것이", 11, False, GRAY_TXT, 4),
    ('계속 마음에 남아, 꿈제작소에서만 플레이할 수 있는', 11, False, GRAY_TXT, 4),
    ('모바일 AR 게임을 기획하게 되었습니다.', 11, False, GRAY_TXT, 14),
    ('아이디어만 전달하기보다 직접 확인하실 수 있는 편이', 11, False, GRAY_TXT, 4),
    ('판단에 도움이 될 것 같아, 실제로 플레이 가능한', 11, False, GRAY_TXT, 4),
    ('형태까지 제작해 두었습니다.', 11, False, GRAY_TXT, 14),
    ('청년의 날 당일 오픈해 방문을 유도하는 용도로,', 11, False, GRAY_TXT, 4),
    ('이후에는 상시 콘텐츠로도 활용할 수 있습니다.', 11, False, GRAY_TXT, 14),
    ('실제 꿈제작소 현장에서 동작을 확인했고,', 11, False, GRAY_TXT, 4),
    ('담당자분께도 시연을 마쳤습니다.', 11, False, GRAY_TXT, 14),
    ('부담 갖지 않고 편하게 검토해 주시면 감사하겠습니다.', 11, False, GRAY_TXT, 4),
]
for i, (txt, sz, bd, col, sa) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = txt
    style_para(p, size=sz, bold=bd, color=col, space_after=sa, line=1.45)

# 오른쪽 — 시연 안내
box(s, 8.35, 1.85, 4.05, 2.9,
    title='시연 영상으로 확인하세요',
    body='실제 꿈제작소에서 촬영한\n플레이 영상입니다.\n\n[ 이 자리에 영상 QR 삽입 ]',
    fill=TEAL_SOFT, line_color=TEAL, title_size=14)

box(s, 8.35, 4.95, 4.05, 1.7,
    title='제안 범위',
    body='청년의 날 행사 사용은\n무상으로 제공합니다.\n\n상시 운영 · 확장은 별도 협의',
    fill=IVORY, line_color=WOOD, title_size=13)

# =====================================================================
# 3. 01 이 게임은?
# =====================================================================
s = new_slide('01', '이 게임은?', page=3)

items = [
    ('01', '실제 위치 기반 AR',
     '포켓몬 GO처럼 GPS로 위치를 인식합니다.\n꿈제작소 안에 있을 때만 카메라가 열리고,\n화면에 실제 공간이 배경으로 나타납니다.'),
    ('02', '청년 캐릭터와 카드 수집',
     '게임을 켠 채로 꿈제작소를 돌아다니면\n층마다 다른 청년 캐릭터가 나타납니다.\n대화에 성공하면 그 청년의 카드를 얻습니다.'),
    ('03', '카드를 모으면 열리는 스토리',
     '카드가 일정 수 이상 모이면\n메인 스토리가 해금됩니다. 실제 꿈제작소\n구조를 옮긴 맵에서 단서를 찾아 나갑니다.'),
    ('04', '미스터리 해결과 배틀',
     '불이 꺼진 꿈제작소에서 단서를 모아\n이상한 박사를 만나고, 모은 카드를\n조합해 배틀로 사건을 해결합니다.'),
]
pos = [(0.9, 1.85), (6.95, 1.85), (0.9, 4.32), (6.95, 4.32)]
for (num, t, b), (x, y) in zip(items, pos):
    rounded(s, x, y, 5.5, 2.28, fill=BOX_BG, radius=0.05)
    add_text(s, x + 0.25, y + 0.42, 1.0, 1.0, num, size=34, bold=True,
             color=RGBColor(0xC9, 0xD4, 0xD9), align=PP_ALIGN.LEFT)
    add_text(s, x + 1.35, y + 0.3, 3.95, 0.4, t, size=14, bold=True,
             color=TEAL_DARK, align=PP_ALIGN.LEFT)
    add_text(s, x + 1.35, y + 0.78, 3.95, 1.3, b, size=10, color=GRAY_TXT,
             align=PP_ALIGN.LEFT, line=1.5)

# =====================================================================
# 4. 01 기대 효과
# =====================================================================
s = new_slide('01', '기대 효과', page=4)

items = [
    ('01', '현장 방문 유도',
     'GPS 반경 안에서만 플레이할 수 있게\n설계했습니다. 게임을 하려면 직접 방문해야\n하므로, 방문 자체가 게임의 조건이 됩니다.'),
    ('02', '청년과의 공감 경험',
     '캐릭터마다 다른 고민을 이야기하고,\n참가자는 대답을 고릅니다. 공감하는 선택을\n했을 때 대화가 이어지도록 설계했습니다.'),
    ('03', '공간 인지도 상승',
     '캐릭터가 카페 · 북스텝 · 스터디룸 등\n실제 공간에 배치되어 있어, 게임을 하는 동안\n자연스럽게 꿈제작소를 알게 됩니다.'),
    ('04', '자발적 SNS 확산',
     '만난 청년 캐릭터를 실제 배경과 함께\n이미지로 저장 · 공유할 수 있습니다. 참가자가\n직접 올리며 추가 홍보비 없이 확산됩니다.'),
]
for (num, t, b), (x, y) in zip(items, pos):
    rounded(s, x, y, 5.5, 2.28, fill=BOX_BG, radius=0.05)
    add_text(s, x + 0.25, y + 0.42, 1.0, 1.0, num, size=34, bold=True,
             color=RGBColor(0xC9, 0xD4, 0xD9), align=PP_ALIGN.LEFT)
    add_text(s, x + 1.35, y + 0.3, 3.95, 0.4, t, size=14, bold=True,
             color=TEAL_DARK, align=PP_ALIGN.LEFT)
    add_text(s, x + 1.35, y + 0.78, 3.95, 1.3, b, size=10, color=GRAY_TXT,
             align=PP_ALIGN.LEFT, line=1.5)

# =====================================================================
# 5. 02 게임의 구조
# =====================================================================
s = new_slide('02', '게임의 구조', page=5)

rows = [
    ('AR 게임 (메인)',
     '꿈제작소에 도착해 게임을 열면 카메라 배경 위에 청년 캐릭터가 나타납니다.\n기기 방향을 돌리면 캐릭터의 화면상 위치도 함께 움직입니다.'),
    ('청년 카드 도감',
     '수집한 카드와 만났던 캐릭터의 대화 기록, 호감도를 다시 확인할 수 있습니다.'),
    ('메인 스토리',
     '카드 레벨이 일정 수치 이상이면 해금됩니다.\n스토리는 위치 제한 없이 어디서든 진행할 수 있습니다.'),
    ('배틀',
     '메인 게임은 배틀이 아닌 만남과 대화 중심입니다.\n배틀은 스토리 진행 중 수집한 카드를 조합해 진행합니다.'),
]
y = 1.85
for label, desc in rows:
    rounded(s, 0.9, y, 2.85, 1.12, fill=YELLOW, line_color=YELLOW_D, radius=0.08)
    tf = add_tb(s, 0.95, y, 2.75, 1.12, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.text = label
    style_para(p, size=13.5, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

    rounded(s, 4.0, y, 8.4, 1.12, fill=BOX_BG, radius=0.05)
    tf = add_tb(s, 4.32, y, 7.8, 1.12, anchor=MSO_ANCHOR.MIDDLE)
    for i, ln in enumerate(desc.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        style_para(p, size=10.5, color=GRAY_TXT, space_after=2, line=1.4)
    y += 1.26

# =====================================================================
# 6. 03 화면 — AR 메인 게임
# =====================================================================
s = new_slide('03', '화면 — AR 메인 게임', page=6)

phone_slot(s, 0.9, 1.85, '① 진입 화면',
           '"꿈제작소 현장에서\n카메라를 열어주세요"')
phone_slot(s, 3.12, 1.85, '② 방향 안내',
           '"왼쪽으로 방향을\n돌려 보세요"\n청년까지 약 13m')
phone_slot(s, 5.34, 1.85, '③ 청년 조우',
           '건물 앞에서 만난\n이도윤\n"말을 걸어볼까요?"')

box(s, 7.75, 1.85, 4.65, 2.16,
    title='게임 방법',
    body='꿈제작소를 돌아다니면서 청년 AR 캐릭터를\n만나고, 대화에 성공해야 카드를 수집합니다.\n\nGPS로 위치를 확인하므로 현장에 있어야\n카메라가 열립니다.',
    title_size=14, align=PP_ALIGN.LEFT)

box(s, 7.75, 4.19, 4.65, 2.16,
    title='캐릭터를 만나지 못하면',
    body='화면에 캐릭터가 보이지 않을 때는 어느 방향으로\n돌아야 하는지 화살표와 남은 거리로 안내합니다.\n\n기기 방향을 돌리면 캐릭터의 화면상 위치도\n함께 움직입니다.',
    title_size=14, align=PP_ALIGN.LEFT)

# =====================================================================
# 7. 03 화면 — 대화 진행
# =====================================================================
s = new_slide('03', '화면 — 청년과의 대화', page=7)

phone_slot(s, 0.9, 1.85, '① 선택지 대화',
           '서고민\n"지금 가장 오래\n붙든 생각은 뭐야?"')
phone_slot(s, 3.12, 1.85, '② 카드 획득',
           '대화 성공\nCARD ARCHIVED\n"서고민 카드 획득"')

box(s, 5.5, 1.85, 6.9, 1.4,
    title='캐릭터를 터치하면 대화가 시작됩니다',
    body='청년마다 성격이 달라 대화 내용이 다르며, 이에 적절한 대답을 고르게 됩니다.',
    fill=TEAL_SOFT, line_color=TEAL, title_size=13.5)

box(s, 5.5, 3.4, 3.35, 1.5,
    title='대화 성공 시',
    body='적합한 대답을 하면\n해당 청년의 카드가\n수집됩니다.',
    fill=RGBColor(0xEF, 0xFA, 0xF2), line_color=RGBColor(0xB8, 0xE5, 0xC8),
    title_size=13, title_color=GREEN_S, body_size=9.5)

box(s, 9.05, 3.4, 3.35, 1.5,
    title='대화 실패 시',
    body='어색한 대답을 고르면\n캐릭터가 사라지고\n대화는 실패합니다.',
    fill=RGBColor(0xFD, 0xF2, 0xF0), line_color=RGBColor(0xF0, 0xC8, 0xC2),
    title_size=13, title_color=RED_SOFT, body_size=9.5)

box(s, 5.5, 5.05, 6.9, 1.3,
    title='실패해도 게임은 계속됩니다',
    body='실패는 종료가 아닙니다. 다시 공간을 돌아다니면\n새로운 청년 신호가 나타나며, 같은 캐릭터를\n다시 만날 수도 있습니다.',
    fill=IVORY, line_color=WOOD, title_size=13, body_size=9.5)

# =====================================================================
# 8. 03 화면 — 청년 캐릭터
# =====================================================================
s = new_slide('03', '화면 — 청년 캐릭터', page=8)

three = [
    ('각종 청년 캐릭터',
     '일하는 청년, 졸린 청년, 꿈꾸는 청년,\n고민하는 청년, 도전하는 청년 등\n여러 모습의 청년을 만날 수 있습니다.'),
    ('청년 캐릭터 속성',
     '캐릭터마다 속성이 있고,\n속성끼리 서로 상성이 있습니다.\n(포켓몬스터의 물 · 불 · 전기 · 풀처럼)'),
    ('카드를 모으면',
     '카드 레벨이 올라가면서\n메인 스토리가 열리고, 모은 카드를\n배틀에서 사용할 수 있습니다.'),
]
for i, (t, b) in enumerate(three):
    box(s, 0.9 + i * 3.92, 1.8, 3.66, 1.85, title=t, body=b, title_size=13)

# 캐릭터 10명
add_text(s, 0.9, 3.95, 11.5, 0.35,
         '등장 청년 캐릭터 — 10명', size=12, bold=True, color=TEAL_DARK)

chars = [
    ('서하늘', '영감'), ('이도윤', '분석'), ('잠유진', '체력'),
    ('서고민', '감정'), ('최도전', '도전'), ('문새봄', '분석'),
    ('박별', '영감'), ('정이든', '체력'), ('강모아', '체력'),
    ('오물결', '감정'),
]
attr_col = {
    '감정': RGBColor(0xFD, 0xE8, 0xE8), '분석': RGBColor(0xE6, 0xF2, 0xFC),
    '도전': RGBColor(0xFE, 0xF2, 0xE0), '체력': RGBColor(0xE6, 0xFA, 0xEE),
    '영감': RGBColor(0xF3, 0xE8, 0xFC),
}
bw, bh = 2.22, 1.05
for i, (nm, at) in enumerate(chars):
    col = i % 5
    row = i // 5
    x = 0.9 + col * 2.33
    y = 4.38 + row * 1.17
    rounded(s, x, y, bw, bh, fill=attr_col[at],
            line_color=RGBColor(0xE2, 0xE2, 0xE2), radius=0.08)
    tf = add_tb(s, x + 0.1, y, bw - 0.2, bh, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.text = nm
    style_para(p, size=12.5, bold=True, color=TEAL_DARK,
               align=PP_ALIGN.CENTER, space_after=1)
    p2 = tf.add_paragraph()
    p2.text = at
    style_para(p2, size=9.5, color=GRAY_MID, align=PP_ALIGN.CENTER)

add_text(s, 0.9, 6.64, 11.5, 0.26,
         '※ 각 캐릭터의 실제 이미지는 게임 화면 캡처로 교체해 주세요.',
         size=8.5, color=RGBColor(0x99, 0x99, 0x99))

# =====================================================================
# 9. 03 화면 — 속성 구조 (오각형)
# =====================================================================
s = new_slide('03', '화면 — 청년 도감')

phone_slot(s, 0.9, 1.85, '청년 도감 상세',
           '이도윤\n호감도 52%\n해금된 대화 3개')

box(s, 3.35, 1.85, 4.4, 2.4,
    title='관계 기록으로 남습니다',
    body='만난 청년마다 호감도가 쌓이고,\n최근 만난 장소가 기록됩니다.\n\n단순한 수집 목록이 아니라\n관계의 기록으로 읽히도록\n설계했습니다.',
    fill=TEAL_SOFT, line_color=TEAL, title_size=13.5, align=PP_ALIGN.LEFT)

box(s, 3.35, 4.43, 4.4, 1.92,
    title='대화가 해금됩니다',
    body='대화를 성공할수록 그 청년과\n나눈 이야기가 하나씩 열립니다.\n\n3턴을 모두 성공하면 전체\n기록을 볼 수 있습니다.',
    title_size=13.5, align=PP_ALIGN.LEFT)

box(s, 8.0, 1.85, 4.4, 2.4,
    title='배틀 역할 안내',
    body='청년마다 속성과 팀 안에서의\n역할이 표시됩니다.\n\n어떤 청년과 조합하면 좋은지\n도감에서 확인할 수 있습니다.',
    title_size=13.5, align=PP_ALIGN.LEFT)

box(s, 8.0, 4.43, 4.4, 1.92,
    title='다시 만나러 가기',
    body='도감에서 청년을 선택하면\n현장 AR로 이동해 같은 청년을\n다시 찾을 수 있습니다.',
    fill=IVORY, line_color=WOOD, title_size=13.5, align=PP_ALIGN.LEFT)

# =====================================================================
# 03 화면 — 청년 속성 구조
# =====================================================================
s = new_slide('03', '화면 — 청년 속성 구조', page=9)

cx, cy, r, d = 3.75, 4.35, 1.78, 1.48
# ★ 배치 순서 = 상성 순서. 시계방향으로 돌면 그대로 규칙이 된다.
labels = ['감정', '분석', '체력', '영감', '도전']
angles = [-90, -18, 54, 126, 198]
centers = []
for lab, ang in zip(labels, angles):
    rad = math.radians(ang)
    centers.append((cx + r * math.cos(rad), cy + r * math.sin(rad)))

# 화살표를 먼저 그려 원 아래에 깔리게 한다
for i in range(5):
    edge_arrow(s, centers[i], centers[(i + 1) % 5], d)

for lab, (ox, oy) in zip(labels, centers):
    oval(s, ox, oy, d, fill=YELLOW, line_color=YELLOW_D)
    tf = add_tb(s, ox - d / 2, oy - d / 2, d, d, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.text = lab
    style_para(p, size=13.5, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

# '이긴다' 라벨 — 오각형 바깥쪽으로 밀어내 배치
for idx in (0, 2):
    a = centers[idx]
    b = centers[(idx + 1) % 5]
    mx, my = (a[0] + b[0]) / 2, (a[1] + b[1]) / 2
    ox, oy = mx - cx, my - cy
    dist = math.hypot(ox, oy) or 1.0
    push = 0.46
    lx = mx + ox / dist * push
    ly = my + oy / dist * push
    add_text(s, lx - 0.42, ly - 0.13, 0.84, 0.28, '이긴다', size=10,
             bold=True, color=RGBColor(0xE8, 0x4C, 0x3D),
             align=PP_ALIGN.CENTER)

tf = add_tb(s, cx - 0.85, cy - 0.45, 1.7, 0.9, anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]
p.text = '청년 속성'
style_para(p, size=11.5, bold=True, color=GRAY_MID, align=PP_ALIGN.CENTER,
           space_after=1)
p2 = tf.add_paragraph()
p2.text = '(강점)'
style_para(p2, size=10, color=GRAY_MID, align=PP_ALIGN.CENTER)

box(s, 7.2, 1.8, 5.2, 1.55,
    title='속성 = 청년의 강점',
    body='청년마다 개인의 강점이 다릅니다.\n게임을 하면서 자신에게 어떤 강점이 있고\n고민을 어떻게 풀어갈지 생각하게 됩니다.',
    fill=TEAL_SOFT, line_color=TEAL, title_size=14)

box(s, 7.2, 3.55, 5.2, 1.55,
    title='상성의 강함과 약함',
    body='배틀 시 속성은 가위바위보처럼\n서로 강하고 약한 관계를 가집니다.\n다만 어떤 속성도 절대적으로 강하지 않습니다.',
    title_size=14)

box(s, 7.2, 5.3, 5.2, 1.3,
    title='왜 상성 구조인가?',
    body="'누가 더 우월한가'가 아니라 '어떤 상황에\n어떤 강점이 필요한가'를 경험하도록 설계했습니다.",
    fill=IVORY, line_color=WOOD, title_size=13)

# 순환 표기
add_text(s, 0.75, 6.62, 6.0, 0.32,
         '감정 → 분석 → 체력 → 영감 → 도전 → 감정',
         size=11.5, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

# =====================================================================
# 10. 03 화면 — 상성 상세
# =====================================================================
s = new_slide('03', '화면 — 속성 상성 규칙', page=10)

pairs = [
    ('감정 → 분석', '분석으로도 설명이 안 되는 감정은 본능의 결과값이다.'),
    ('분석 → 체력', '정확한 분석은 체력으로 버틴 것보다 더 결론을 낸다.'),
    ('체력 → 영감', '영감이 있어도 체력이 부족하면 실현하기 어렵다.'),
    ('영감 → 도전', '도전할 용기가 있어도 영감 없이는 행동일 뿐이다.'),
    ('도전 → 감정', '감정으로 포기한 기회를 도전으로 이겨낸다.'),
]
y = 1.85
for head, desc in pairs:
    rounded(s, 0.9, y, 2.6, 0.86, fill=YELLOW, line_color=YELLOW_D, radius=0.1)
    tf = add_tb(s, 0.95, y, 2.5, 0.86, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.text = head
    style_para(p, size=12.5, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

    rounded(s, 3.72, y, 8.68, 0.86, fill=BOX_BG, radius=0.06)
    tf = add_tb(s, 4.05, y, 8.1, 0.86, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.text = desc
    style_para(p, size=10.5, color=GRAY_TXT)
    y += 0.99

box(s, 0.9, 6.85, 11.5, 0.0)  # spacer (무시)

# =====================================================================
# 11. 03 화면 — 배틀
# =====================================================================
s = new_slide('03', '화면 — 메인 스토리 배틀', page=11)

box(s, 0.9, 1.8, 5.55, 1.5,
    title='배틀 게임 방식',
    body='수집한 카드 중 3장을 골라 카드 칸에 넣습니다.\n상대의 카드와 만나 속성에 따라 승패가 갈립니다.',
    fill=TEAL_SOFT, line_color=TEAL, title_size=14)

box(s, 0.9, 3.5, 5.55, 1.62,
    title='카드 조합 종류',
    body='페어 — 2장이 같은 속성\n무지개 — 3장이 모두 다른 속성\n몰빵 — 3장이 모두 같은 속성',
    title_size=14)

box(s, 0.9, 5.32, 5.55, 1.62,
    title='대표 속성 결정 방식',
    body='페어 — 겹친 2장의 속성이 대표\n무지개 — 가장 먼저 낸 카드가 대표\n몰빵 — 3장의 공통 속성이 대표',
    title_size=14)

box(s, 6.85, 1.8, 5.55, 2.4,
    title='승패 판정',
    body='카드 조합으로 결정된 대표 속성으로\n나와 상대의 승패가 결정됩니다.\n\n감정 → 분석  ·  분석 → 체력  ·  체력 → 영감\n영감 → 도전  ·  도전 → 감정\n\n화살표 앞 속성이 이깁니다.',
    title_size=14)

box(s, 6.85, 4.4, 3.35, 2.54,
    title='최종 보스',
    body='이상한 박사는 직전 턴에\n플레이어가 낸 대표 속성을\n다음 턴에 복제합니다.\n\n매 턴 조합을 바꾸는\n전략이 필요합니다.\n\n3라운드 중 2승하면\n승리합니다.',
    fill=IVORY, line_color=WOOD, title_size=13, body_size=9.5)

phone_slot(s, 10.32, 1.85, '이상한 박사 배틀',
           '체력 게이지\n감정 vs 분석\n카드 3장 선택')

# =====================================================================
# 12. 03 화면 — 메인 스토리 진행
# =====================================================================
s = new_slide('03', '화면 — 메인 스토리 진행', page=12)

phone_slot(s, 0.9, 1.85, '① 야간 탐사 맵',
           '코워킹A · 북스텝\n카페테리아\n바닥을 터치해 이동')
phone_slot(s, 3.12, 1.85, '② 한숨덩이 조우',
           '"한숨덩이가\n길을 막았다"\n우회 / 배틀 선택')

box(s, 5.5, 1.85, 6.9, 1.5,
    title='메인 스토리가 열리면',
    body='실제 꿈제작소의 건물 구조로 맵이 제작됩니다. 이 맵을 이동하며\n단서를 모으고, 최종 보스와 배틀 후 결말을 봅니다.',
    fill=TEAL_SOFT, line_color=TEAL, title_size=13.5, body_size=9.5)

box(s, 5.5, 3.5, 6.9, 1.28,
    title='조작 방식',
    body='화면을 터치한 지점으로 이동하며, 실제 위치 기반이 아닙니다.\n위치 제한 없이 어디서든 진행할 수 있습니다.',
    title_size=13.5, body_size=9.5)

box(s, 5.5, 4.93, 6.9, 1.42,
    title='플레이 흐름',
    body='스토리 해금 → 불이 꺼진 꿈제작소 탐색 → 증거 수집\n→ 화가를 만나 대화 → 한숨덩이 배틀로 상처약 습득\n→ 이상한 박사와 대화 후 배틀 → 선택지에 따라 다른 엔딩',
    fill=IVORY, line_color=WOOD, title_size=13.5, body_size=9.5)

# =====================================================================
# 03 화면 — 스토리 맵 구성
# =====================================================================
s = new_slide('03', '화면 — 스토리 맵 구성')

box(s, 0.9, 1.85, 7.3, 1.5,
    title='실제 건물 구조를 그대로 옮겼습니다',
    body='참가자가 평소 이용하는 공간이 그대로 게임 맵이 되기 때문에,\n스토리를 진행하는 동안 자연스럽게 꿈제작소를 익히게 됩니다.',
    fill=TEAL_SOFT, line_color=TEAL, title_size=13.5, body_size=9.5)

box(s, 0.9, 3.5, 7.3, 1.4,
    title='층간 이동',
    body='계단과 엘리베이터 표식을 터치하면 다른 층으로 이동합니다.\n벽과 가구에는 막혀서 우회 경로를 찾아야 합니다.',
    title_size=13.5, body_size=9.5)

box(s, 0.9, 5.05, 7.3, 1.3,
    title='단서와 아이템',
    body='층마다 고유한 단서와 아이템이 배치되어 있고,\n획득한 물품은 가방에서 다시 확인할 수 있습니다.',
    fill=IVORY, line_color=WOOD, title_size=13.5, body_size=9.5)

# 맵 구성
rounded(s, 8.35, 1.85, 4.05, 4.5, fill=BOX_BG, radius=0.04)
tf = add_tb(s, 8.6, 2.02, 3.55, 4.2)
map_lines = [
    ('메인 스토리 맵 구성', 13, True, TEAL_DARK, 10, PP_ALIGN.CENTER),
    ('A동 · B동 · 중간 통로로 이동', 9.5, False, GRAY_MID, 12, PP_ALIGN.CENTER),
    ('A동', 11, True, TEAL_DARK, 4, PP_ALIGN.LEFT),
    ('1층  로비(코워킹A), 북스텝,', 9.5, False, GRAY_TXT, 1, PP_ALIGN.LEFT),
    ('        카페, 스튜디오, 스터디룸', 9.5, False, GRAY_TXT, 5, PP_ALIGN.LEFT),
    ('2층  코워킹B, 청년거버넌스룸,', 9.5, False, GRAY_TXT, 1, PP_ALIGN.LEFT),
    ('        컨설팅룸 1~4, 세미나실 1~2', 9.5, False, GRAY_TXT, 5, PP_ALIGN.LEFT),
    ('3층  세미나실 3~4, 클래스룸 1~3', 9.5, False, GRAY_TXT, 5, PP_ALIGN.LEFT),
    ('4층  강당, 인터뷰룸, 라운지', 9.5, False, GRAY_TXT, 12, PP_ALIGN.LEFT),
    ('B동', 11, True, TEAL_DARK, 4, PP_ALIGN.LEFT),
    ('내일꿈갤러리, 신체활동실', 9.5, False, GRAY_TXT, 5, PP_ALIGN.LEFT),
]
for i, (txt, sz, bd, col, sa, al) in enumerate(map_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = txt
    style_para(p, size=sz, bold=bd, color=col, space_after=sa, align=al,
               line=1.3)

# =====================================================================
# 13. 03 화면 — 메인 스토리 내용
# =====================================================================
s = new_slide('03', '화면 — 메인 스토리 내용', page=13)

rounded(s, 0.9, 1.8, 7.0, 5.14, fill=BOX_BG, radius=0.04)
tf = add_tb(s, 1.25, 2.05, 6.35, 4.7)
story = [
    ('메인 스토리', 14, True, TEAL_DARK, 12),
    ('주인공은 꿈제작소의 소문을 듣습니다.', 10.5, False, GRAY_TXT, 6),
    ('"폐관 후 유령이 돌아다니는 것을 봤다.', 10.5, False, GRAY_MID, 1),
    ('아침에 개관하면 화장실 세 번째 칸은 뚜껑이 닫혀 있다.', 10.5, False, GRAY_MID, 1),
    ('청년 작가의 캔버스는 개관하면 뒤집혀 있다."', 10.5, False, GRAY_MID, 10),
    ('주인공은 밤에 꿈제작소로 가서 비밀을 파헤치기로 합니다.', 10.5, False, GRAY_TXT, 6),
    ('그러다 이상한 박사를 만납니다. 박사는 청년의 꿈 카드를', 10.5, False, GRAY_TXT, 1),
    ('연구하며 복제하고 기록을 남기고 있지만, 복제할수록', 10.5, False, GRAY_TXT, 1),
    ('원본 카드의 힘은 약해집니다.', 10.5, False, GRAY_TXT, 10),
    ('스토리 중간에 화가를 만나고, 그림이 뒤집어진 이유를', 10.5, False, GRAY_TXT, 1),
    ('알게 됩니다. 건물 곳곳의 한숨덩이와 마주치면 배틀이', 10.5, False, GRAY_TXT, 1),
    ('시작됩니다. 박사와의 대화 끝에 최종 배틀로 이어집니다.', 10.5, False, GRAY_TXT, 10),
]
for i, (txt, sz, bd, col, sa) in enumerate(story):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = txt
    style_para(p, size=sz, bold=bd, color=col, space_after=sa, line=1.45)

box(s, 8.35, 1.8, 4.05, 2.45,
    title='이상한 박사?',
    body='청년들의 꿈 카드를 몰래 복제하며\n기록을 남기는 인물입니다.\n복제할수록 원본 카드의 힘이\n약해지지만 멈추지 않습니다.',
    fill=IVORY, line_color=WOOD, title_size=14)

box(s, 8.35, 4.45, 4.05, 2.49,
    title='두 가지 엔딩',
    body='배틀 후 선택지에 따라\n결말이 갈립니다.\n\n① 복제 장치를 멈추고\n     원본을 지킨다\n\n② 박사를 설득해\n     기록을 함께 보존한다',
    fill=TEAL_SOFT, line_color=TEAL, title_size=14)

# =====================================================================
# 14. 04 운영 안내 ① 개인정보
# =====================================================================
s = new_slide('04', '기관 운영 안내 ① 개인정보', page=14)

rounded(s, 0.9, 1.75, 11.5, 0.78, fill=TEAL_DARK, line_color=None, radius=0.12)
tf = add_tb(s, 1.0, 1.75, 11.3, 0.78, anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]
p.text = '이 게임은 참가자의 개인정보를 수집하지 않습니다.'
style_para(p, size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box(s, 0.9, 2.75, 5.6, 2.05,
    title='수집하지 않는 것',
    body='· 이름 · 전화번호 · 이메일\n· 카메라 영상 (서버로 전송되지 않음)\n· GPS 이동 이력\n· 로그인 · 회원가입 정보',
    fill=RGBColor(0xFD, 0xF2, 0xF0), line_color=RGBColor(0xF0, 0xC8, 0xC2),
    title_size=14, title_color=RED_SOFT, align=PP_ALIGN.LEFT)

box(s, 6.8, 2.75, 5.6, 2.05,
    title='대신 이렇게 처리합니다',
    body='· 게임 기록은 참가자 휴대폰에만 저장\n· 참가자가 직접 삭제하는 버튼 제공\n· 첫 실행 시 저장 방식 안내 팝업\n· 경품은 익명 1회용 코드로만 처리',
    fill=RGBColor(0xEF, 0xFA, 0xF2), line_color=RGBColor(0xB8, 0xE5, 0xC8),
    title_size=14, title_color=GREEN_S, align=PP_ALIGN.LEFT)

add_text(s, 0.9, 5.0, 11.5, 0.35, '기관 입장에서의 이점',
         size=13, bold=True, color=TEAL_DARK)

merits = [
    ('개인정보 동의서\n불필요', '수집 항목이 없어 별도 동의\n절차가 발생하지 않습니다.'),
    ('데이터 보관 · 파기\n업무 없음', '행사 후 정리할 개인정보가\n존재하지 않습니다.'),
    ('유출 위험\n없음', '서버에 개인식별정보가\n저장되지 않습니다.'),
]
for i, (t, b) in enumerate(merits):
    box(s, 0.9 + i * 3.92, 5.42, 3.66, 1.5, title=t, body=b,
        fill=IVORY, line_color=WOOD, title_size=12.5, body_size=9.5)

# =====================================================================
# 15. 04 운영 안내 ② 준비물
# =====================================================================
s = new_slide('04', '기관 운영 안내 ② 필요한 준비물', page=15)

rows = [
    ('참가자', '스마트폰과 인터넷만 있으면 됩니다. 앱 설치 · 회원가입 · 전화번호 입력이 없어\nQR 스캔 후 1분 내에 시작할 수 있습니다.'),
    ('운영 인력', '안내 1명, 경품 데스크 1명이면 충분합니다.\n완주 화면을 확인하고 코드 발급 버튼을 누르는 것이 전부입니다.'),
    ('비치물', 'QR이 인쇄된 안내물 또는 배너, 무작위 번호가 적힌 종이 입장권, 예비 수령 대장'),
    ('소요 시간', 'AR 체험만 약 15분, 메인 스토리까지 완주하면 30~60분'),
    ('사전 준비', 'QR 안내물 출력과 종이 입장권 준비만 하면 됩니다.\n게임은 추가 개발 없이 현재 상태로 바로 사용할 수 있습니다.'),
    ('개발 비용', '이미 제작이 완료된 상태로 제공합니다. 경품비만 기관에서 결정하시면 됩니다.'),
]
y = 1.8
for label, desc in rows:
    rounded(s, 0.9, y, 2.3, 0.78, fill=YELLOW, line_color=YELLOW_D, radius=0.1)
    tf = add_tb(s, 0.95, y, 2.2, 0.78, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.text = label
    style_para(p, size=12.5, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

    rounded(s, 3.42, y, 8.98, 0.78, fill=BOX_BG, radius=0.06)
    tf = add_tb(s, 3.72, y, 8.4, 0.78, anchor=MSO_ANCHOR.MIDDLE)
    for i, ln in enumerate(desc.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        style_para(p, size=10, color=GRAY_TXT, space_after=1, line=1.32)
    y += 0.86

# =====================================================================
# 16. 04 운영 안내 ③ 행사 부스 운영  (신규)
# =====================================================================
s = new_slide('04', '기관 운영 안내 ③ 행사 부스 운영', page=16)

flow = [
    ('01  안내', '부스에서 QR을 스캔하고\n게임 방법을 30초 안내받습니다.\n종이 입장권을 함께 받습니다.'),
    ('02  체험 시작', '건물 안으로 이동해\n카메라를 열면 청년 캐릭터가\n나타납니다.'),
    ('03  자유 탐험', '층을 옮겨 다니며 청년을 만나고\n카드를 모읍니다. 부스로 돌아올\n필요가 없습니다.'),
    ('04  경품 수령', '완주 화면과 입장권을 제시하면\n운영자가 코드를 발급하고\n경품을 전달합니다.'),
]
for i, (t, b) in enumerate(flow):
    x = 0.9 + i * 2.93
    rounded(s, x, 1.8, 2.68, 1.42, fill=BOX_BG, radius=0.05)
    tf = add_tb(s, x + 0.18, 1.92, 2.32, 1.2)
    p = tf.paragraphs[0]
    p.text = t
    style_para(p, size=12.5, bold=True, color=TEAL_DARK, space_after=5)
    for ln in b.split('\n'):
        p = tf.add_paragraph()
        p.text = ln
        style_para(p, size=9, color=GRAY_TXT, space_after=1, line=1.3)
    if i < 3:
        add_text(s, x + 2.68, 2.36, 0.25, 0.3, '›', size=18, bold=True,
                 color=TEAL, align=PP_ALIGN.CENTER)

booth = [
    ('입구 안내 부스', '담당 1명 · QR 안내와 게임 방법 설명, 종이 입장권 배부.\n참가자를 붙잡아 둘 필요가 없어 회전이 빠릅니다.'),
    ('경품 데스크', '담당 1명 · 완주 화면과 입장권 번호를 확인한 뒤 코드를 발급합니다.\n1건 처리에 약 20초가 걸립니다.'),
    ('동시 수용 인원', '참가자가 각자 휴대폰으로 플레이하므로 인원 제한이 없습니다.\n대기가 생길 수 있는 지점은 경품 데스크 한 곳뿐입니다.'),
]
y = 3.4
for label, desc in booth:
    rounded(s, 0.9, y, 2.6, 0.76, fill=TEAL_SOFT, line_color=TEAL, radius=0.1)
    tf = add_tb(s, 0.95, y, 2.5, 0.76, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.text = label
    style_para(p, size=12, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

    rounded(s, 3.72, y, 8.68, 0.76, fill=BOX_BG, radius=0.06)
    tf = add_tb(s, 4.0, y, 8.2, 0.76, anchor=MSO_ANCHOR.MIDDLE)
    for i, ln in enumerate(desc.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        style_para(p, size=9.5, color=GRAY_TXT, space_after=1, line=1.3)
    y += 0.84

box(s, 0.9, 6.0, 5.6, 0.9,
    title='네트워크 장애 대응',
    body='인터넷이 끊겨도 이미 열린 게임은 계속 진행됩니다.\n경품은 종이 입장권과 확인 도장으로 처리할 수 있습니다.',
    fill=IVORY, line_color=WOOD, title_size=12, body_size=9.5)

box(s, 6.8, 6.0, 5.6, 0.9,
    title='접근성',
    body='카메라 · GPS 권한 사용이 어려운 참가자를 위해\n위치 제한 없는 체험 모드를 제공합니다.',
    fill=IVORY, line_color=WOOD, title_size=12, body_size=9.5)

# =====================================================================
# 16. 05 활용 흐름
# =====================================================================
s = new_slide('05', '활용 흐름', page=17)

steps = [
    ('STEP 01', 'SNS · 현장 홍보',
     '청년의 날에 맞춰 게임을\n오픈하고, 포스터 · SNS ·\n현장 배너에 QR을 배치합니다.'),
    ('STEP 02', '방문 후 게임 시작',
     '앱 설치와 회원가입이 없어\nQR 하나로 바로 시작합니다.\n진입 장벽이 거의 없습니다.'),
    ('STEP 03', '공간을 돌며 체험',
     '층마다 다른 청년을 만나며\n카페 · 북스텝 · 스터디룸 등\n실제 공간을 알게 됩니다.'),
    ('STEP 04', '완주와 재확산',
     '스토리를 완주하고 경품을\n받습니다. 기록 이미지를 SNS에\n공유하며 홍보가 이어집니다.'),
]
bw = 2.72
for i, (st, t, b) in enumerate(steps):
    x = 0.9 + i * 2.95
    shade = [RGBColor(0xF9, 0xF9, 0xF9), RGBColor(0xF2, 0xF5, 0xF6),
             RGBColor(0xE9, 0xEF, 0xF1), RGBColor(0xDE, 0xE8, 0xEB)][i]
    rounded(s, x, 1.9, bw, 0.72, fill=shade, line_color=BORDER, radius=0.3)
    tf = add_tb(s, x, 1.9, bw, 0.72, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.text = st
    style_para(p, size=12.5, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

    if i < 3:
        oval(s, x + bw + 0.115, 2.26, 0.5, fill=YELLOW, line_color=None)
        tf = add_tb(s, x + bw - 0.135, 2.26 - 0.25, 0.5, 0.5,
                    anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.text = '▶'
        style_para(p, size=10, bold=True, color=TEAL_DARK,
                   align=PP_ALIGN.CENTER)

    box(s, x, 2.82, bw, 2.15, title=t, body=b, title_size=12.5,
        body_size=9.5)

rounded(s, 0.9, 5.3, 11.5, 0.85, fill=YELLOW, line_color=YELLOW_D, radius=0.1)
tf = add_tb(s, 1.0, 5.3, 11.3, 0.85, anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]
p.text = '홍보 → 방문 → 체험 → 재확산이 하나의 흐름으로 이어집니다.'
style_para(p, size=14, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

box(s, 0.9, 6.32, 11.5, 0.62,
    body='행사 이후에는 신규 방문자를 위한 상시 공간 안내 콘텐츠로 전환해 활용할 수 있습니다.',
    body_size=10.5, fill=BOX_BG)

# =====================================================================
# 17. 06 진행 제안
# =====================================================================
s = new_slide('06', '진행 제안', page=18)

box(s, 0.9, 1.8, 5.6, 2.5,
    title='완성 및 검증된 범위',
    body='· AR 청년 조우와 대화 · 카드 수집\n· 청년 도감과 호감도 기록\n· 메인 스토리 탐사와 배틀\n· 운영자용 경품 코드 화면\n· 실제 꿈제작소 현장 동작 확인 완료',
    fill=RGBColor(0xEF, 0xFA, 0xF2), line_color=RGBColor(0xB8, 0xE5, 0xC8),
    title_size=14, title_color=GREEN_S, align=PP_ALIGN.LEFT)

box(s, 6.8, 1.8, 5.6, 2.5,
    title='기관에서 결정해 주실 것',
    body='· 청년의 날 사용 여부\n· 경품 종류와 수령 기준\n· 종이 입장권 번호 규칙\n· 안내 · 경품 데스크 담당자 각 1명\n· QR 안내물 또는 배너 제작',
    fill=TEAL_SOFT, line_color=TEAL, title_size=14, align=PP_ALIGN.LEFT)

box(s, 0.9, 4.5, 5.6, 1.42,
    title='안정적인 운영 환경',
    body='만료 없는 상용 호스팅에 배포해 운영합니다.\n주 링크와 예비 링크를 이중으로 준비해\n행사 당일 장애에 대비합니다.',
    fill=TEAL_SOFT, line_color=TEAL, title_size=13, body_size=9.5)

rounded(s, 6.8, 4.5, 5.6, 1.42, fill=IVORY, line_color=WOOD, radius=0.05)
tf = add_tb(s, 7.1, 4.66, 5.0, 1.1)
prov = [
    ('제공 범위', 13, True, TEAL_DARK, 5, PP_ALIGN.CENTER),
    ('행사 사용은 무상으로 제공합니다.', 9.5, False, GRAY_TXT, 2, PP_ALIGN.LEFT),
    ('상시 운영 · 콘텐츠 추가 · 호스팅 이전은', 9.5, False, GRAY_TXT, 2, PP_ALIGN.LEFT),
    ('별도 협의가 필요합니다.', 9.5, False, GRAY_TXT, 2, PP_ALIGN.LEFT),
    ('기획과 저작권은 제작자에게 있습니다.', 9.5, False, GRAY_TXT, 0, PP_ALIGN.LEFT),
]
for i, (txt, sz, bd, col, sa, al) in enumerate(prov):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = txt
    style_para(p, size=sz, bold=bd, color=col, space_after=sa, align=al,
               line=1.35)

rounded(s, 0.9, 6.1, 11.5, 0.84, fill=TEAL_DARK, line_color=None, radius=0.1)
tf = add_tb(s, 1.2, 6.1, 11.1, 0.84, anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]
p.text = '현장 테스트는 이미 완료했습니다.'
style_para(p, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
           space_after=2)
p2 = tf.add_paragraph()
p2.text = '진행 여부를 알려주시면 행사 일정에 맞춰 준비하겠습니다.'
style_para(p2, size=10.5, color=RGBColor(0xB8, 0xD8, 0xDF),
           align=PP_ALIGN.CENTER)

# =====================================================================
out = '/projects/sandbox/pptx-build/꿈제작소_AR게임_제안서_이아선.pptx'
prs.save(out)
print('저장 완료:', out)
print('총 슬라이드 수:', len(prs.slides.__iter__.__self__._sldIdLst))
